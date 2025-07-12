from pathlib import Path
import os
from langchain_community.document_loaders.pdf import PyPDFLoader
from langchain_community.document_loaders import Docx2txtLoader, TextLoader
from langchain.docstore.document import Document as LangchainDocument 
from pptx import Presentation 

from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_openai.embeddings import OpenAIEmbeddings
from langchain_community.vectorstores.faiss import FAISS
from langchain.chains.conversational_retrieval.base import ConversationalRetrievalChain
from langchain_openai.chat_models import ChatOpenAI
from langchain.memory import ConversationBufferMemory
import streamlit as st
from dotenv import load_dotenv, find_dotenv

_ = load_dotenv(find_dotenv())

pasta_arquivos = Path(__file__).parent / "arquivos"
pasta_arquivos.mkdir(exist_ok=True)

nome_modelo = "gpt-3.5-turbo-0125"

def importar_documentos():
    
    documentos = []
    
    
    for arquivo in pasta_arquivos.glob("*.pdf"):
        try:
            loader = PyPDFLoader(str(arquivo))
            docs_arquivo = loader.load()
            documentos.extend(docs_arquivo)
        except Exception as e:
            st.error(f"Erro ao carregar PDF {arquivo.name}: {e}")

    
    for arquivo in pasta_arquivos.glob("*.docx"):
        try:
            loader = Docx2txtLoader(str(arquivo))
            docs_arquivo = loader.load()
            documentos.extend(docs_arquivo)
        except Exception as e:
            st.error(f"Erro ao carregar DOCX {arquivo.name}: {e}")

    
    for arquivo in pasta_arquivos.glob("*.pptx"):
        try:
            prs = Presentation(arquivo)
            texto_slides = []
            for i, slide in enumerate(prs.slides):
                slide_texto = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_texto.append(shape.text)
                if slide_texto:
                    texto_completo_slide = "\n".join(slide_texto)
                    documentos.append(LangchainDocument(page_content=texto_completo_slide, metadata={"source": arquivo.name, "page": i + 1}))
        except Exception as e:
            st.error(f"Erro ao carregar PPTX {arquivo.name}: {e}")

    
    for arquivo in pasta_arquivos.glob("*.txt"):
        try:
            loader = TextLoader(str(arquivo), encoding='utf-8')
            docs_arquivo = loader.load()
            documentos.extend(docs_arquivo)
        except Exception as e:
            try:
                loader = TextLoader(str(arquivo), encoding='latin-1')
                docs_arquivo = loader.load()
                documentos.extend(docs_arquivo)
            except Exception as e2:
                st.error(f"Erro ao carregar TXT {arquivo.name}: {e2}")

    return documentos

def dividir_documentos(documentos):
    
    divisor = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=50,
        separators=["\n\n", "\n", ".", " ", ""]
    )

    documentos_divididos = divisor.split_documents(documentos)
    for i, doc in enumerate(documentos_divididos):
        if "source" in doc.metadata:
            doc.metadata["source"] = Path(doc.metadata["source"]).name
        doc.metadata["id"] = i

    return documentos_divididos

def criar_vetores(documentos):
    
    modelo_embeddings = OpenAIEmbeddings()
    vetores = FAISS.from_documents(
        documents=documentos,
        embedding=modelo_embeddings
    )

    return vetores

def criar_conversa():
    
    st.sidebar.text("Carregando documentos...")

    documentos = importar_documentos()
    if not documentos:
        st.error("Nenhum documento encontrado ou carregado!")
        return None

    st.sidebar.text("Processando textos...")
    documentos_divididos = dividir_documentos(documentos)

    if not documentos_divididos:
        st.error("Nenhum pedaço de documento foi gerado após a divisão. O arquivo pode estar vazio ou o splitter não conseguiu processá-lo!")
        return None

    st.sidebar.text("Criando embeddings...")
    vetores = criar_vetores(documentos_divididos)

    st.sidebar.text("Configurando modelo...")
    chat = ChatOpenAI(
        model=nome_modelo,
        temperature=0.7
    )

    memoria = ConversationBufferMemory(
        return_messages=True,
        memory_key="chat_history",
        output_key="answer"
    )

    recuperador = vetores.as_retriever(
        search_kwargs={"k": 3}
    )

    chain_conversa = ConversationalRetrievalChain.from_llm(
        llm=chat,
        memory=memoria,
        retriever=recuperador,
        return_source_documents=True
    )
    st.session_state["chain"] = chain_conversa
    return chain_conversa
